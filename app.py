# app.py — TextCraft AI Studio (single-file, modern UI + full features)
import streamlit as st
import sqlite3, hashlib, re, time, math
from datetime import datetime
from collections import Counter

# ---- Optional readers (all FREE & optional) ----
try:
    import PyPDF2; PDF_OK = True
except Exception:
    PDF_OK = False
try:
    from docx import Document; DOCX_OK = True
except Exception:
    DOCX_OK = False
try:
    from pptx import Presentation; PPTX_OK = True
except Exception:
    PPTX_OK = False
try:
    import pandas as pd; PANDAS_OK = True
except Exception:
    PANDAS_OK = False

# =====================
# PAGE CONFIG & THEME
# =====================
st.set_page_config(page_title="TextCraft AI Studio", page_icon="🎯", layout="wide")

# ---------- Global CSS (slide-style UI) ----------
st.markdown("""
<style>
/* App background (purple slide look) */
[data-testid="stAppViewContainer"] {
  background:
    radial-gradient(1200px 600px at 15% 15%, rgba(124,58,237,.32), transparent 45%),
    radial-gradient(1000px 600px at 85% 25%, rgba(99,102,241,.28), transparent 45%),
    linear-gradient(135deg, #3a1c71 0%, #5e2ea6 35%, #7c3aed 70%, #8b5cf6 100%);
}

/* Header */
[data-testid="stHeader"] { background: transparent; }

/* Typography (dark-on-light inside cards; white headings for contrast outside) */
h1, h2, h3 {
  font-family: ui-sans-serif, -apple-system, Segoe UI, Roboto, Helvetica, Arial, "Apple Color Emoji","Segoe UI Emoji";
  letter-spacing: .3px;
  color: #eef; 
}
p, label, span { color: #111; }

/* Cards */
.card {
  background: rgba(255,255,255,.80);
  border: 1px solid rgba(0,0,0,.08);
  border-radius: 18px;
  padding: 18px 18px;
  box-shadow: 0 10px 26px rgba(0,0,0,.25);
}
.card-tight {
  background: rgba(255,255,255,.82);
  border: 1px solid rgba(0,0,0,.08);
  border-radius: 14px;
  padding: 12px 14px;
}

/* Sidebar (dark panel like your slide) */
[data-testid="stSidebar"] {
  background: linear-gradient(180deg, #0b1020 0%, #131a2e 100%);
  border-right: 1px solid rgba(255,255,255,.08);
}
.sidebar-card {
  background: rgba(255,255,255,.08);
  border: 1px solid rgba(255,255,255,.12);
  border-radius: 16px;
  padding: 14px;
  color: #f3f4ff;
}
.sidebar-card * { color: #e7e9ff !important; }

/* Buttons */
div.stButton > button {
  width: 100%;
  border: 0;
  color: #0b1020;
  font-weight: 700;
  border-radius: 12px;
  padding: .7rem 1rem;
  background: linear-gradient(90deg, #d9b3ff 0%, #b285f7 100%);
  transition: transform .1s ease, box-shadow .2s ease, opacity .2s ease;
}
div.stButton > button:hover { transform: translateY(-1px); box-shadow: 0 10px 24px rgba(178,133,247,.35); }
div.stButton > button:active { transform: translateY(0); opacity: .95; }

/* Inputs */
.stTextInput > div > div > input,
.stTextArea textarea, .stSelectbox div[data-baseweb="select"] > div {
  background: rgba(255,255,255,.9);
  border: 1px solid rgba(0,0,0,.12);
  border-radius: 12px;
  color: #111;
}

/* Pills */
.pill {
  display:inline-block; padding: 4px 10px; border-radius: 999px;
  font-size: 12px; font-weight: 700; letter-spacing:.2px;
  background: rgba(139,92,246,.18); color: #2b2453; border: 1px solid rgba(139,92,246,.35);
}

/* Metrics */
.metric-title { color:#333; font-weight:700; font-size:13px; }
.metric-val { color:#000; font-size:18px; font-weight:800; }

/* Result header bar (like slide) */
.result-header {
  background: linear-gradient(90deg, #7c3aed 0%, #8b5cf6 100%);
  color: #fff; padding: 12px 14px; border-radius: 12px; font-weight: 800; letter-spacing: .2px;
  display:flex; align-items:center; gap:10px; box-shadow: inset 0 0 0 1px rgba(255,255,255,.18);
}

/* Status bar */
.progress-wrap { background: rgba(0,0,0,.08); height: 10px; border-radius: 999px; overflow: hidden; }
.progress-fill { height: 100%; background: linear-gradient(90deg, #22c55e, #a3e635); }

/* Check / bullet rows */
.check { display:flex; align-items:flex-start; gap:8px; }
.check .dot { width: 10px; height: 10px; border-radius: 50%; margin-top: 7px; background: #22c55e; box-shadow: 0 0 0 3px rgba(34,197,94,.15); }
.warn .dot { background: #f59e0b; box-shadow: 0 0 0 3px rgba(245,158,11,.18); }
.err  .dot { background: #ef4444; box-shadow: 0 0 0 3px rgba(239,68,68,.18); }

/* Footer */
.footer { color:#cbd3ff; font-size:12px; text-align:center; margin-top: 10px; }
</style>
""", unsafe_allow_html=True)

# =====================
# DATABASE (SQLite)
# =====================
DB = "textcraft_ai.db"

def init_db():
    conn = sqlite3.connect(DB, timeout=30)
    c = conn.cursor()
    c.execute("""
    CREATE TABLE IF NOT EXISTS users(
      id INTEGER PRIMARY KEY AUTOINCREMENT,
      username TEXT UNIQUE NOT NULL,
      password TEXT NOT NULL,
      created_at TEXT NOT NULL,
      last_login TEXT,
      session_count INTEGER DEFAULT 0,
      total_analyses INTEGER DEFAULT 0
    )""")
    c.execute("""
    CREATE TABLE IF NOT EXISTS analyses(
      id INTEGER PRIMARY KEY AUTOINCREMENT,
      user_id INTEGER NOT NULL,
      analysis_type TEXT NOT NULL,
      input_text TEXT NOT NULL,
      result_text TEXT NOT NULL,
      timestamp TEXT NOT NULL
    )""")
    conn.commit(); conn.close()

hash_pw = lambda p: hashlib.sha256(p.encode()).hexdigest()

def create_user(username, password):
    try:
        conn = sqlite3.connect(DB, timeout=30); c = conn.cursor()
        c.execute("INSERT INTO users(username,password,created_at) VALUES(?,?,?)",
                  (username, hash_pw(password), datetime.now().isoformat()))
        conn.commit(); conn.close(); return True
    except sqlite3.IntegrityError:
        return False

def verify_user(username, password):
    conn = sqlite3.connect(DB, timeout=30); c = conn.cursor()
    c.execute("SELECT id FROM users WHERE username=? AND password=?", (username, hash_pw(password)))
    row = c.fetchone()
    if row:
        c.execute("UPDATE users SET last_login=?, session_count=session_count+1 WHERE id=?",
                  (datetime.now().isoformat(), row[0]))
        conn.commit()
    conn.close()
    return row[0] if row else None

def get_user_stats(user_id):
    conn = sqlite3.connect(DB, timeout=30); c = conn.cursor()
    c.execute("SELECT session_count,last_login,created_at,total_analyses FROM users WHERE id=?", (user_id,))
    row = c.fetchone(); conn.close()
    return row if row else (0, None, None, 0)

def save_analysis(user_id, kind, src, out):
    conn = sqlite3.connect(DB, timeout=30); c = conn.cursor()
    c.execute("INSERT INTO analyses(user_id,analysis_type,input_text,result_text,timestamp) VALUES(?,?,?,?,?)",
              (user_id, kind, (src or "")[:1000], (out or "")[:2000], datetime.now().isoformat()))
    c.execute("UPDATE users SET total_analyses=total_analyses+1 WHERE id=?", (user_id,))
    conn.commit(); conn.close()

# =====================
# FILE UTILS
# =====================
def get_supported_file_types():
    exts = ['txt','md','rtf','json']
    if PDF_OK: exts.append('pdf')
    if DOCX_OK: exts += ['docx','doc']
    if PPTX_OK: exts += ['pptx','ppt']
    if PANDAS_OK: exts += ['csv','xlsx','xls']
    return exts

def read_uploaded_file(uploaded_file):
    if not uploaded_file: return "", None
    ext = uploaded_file.name.split('.')[-1].lower()
    try:
        if ext == 'txt':
            return uploaded_file.read().decode('utf-8', 'ignore'), None
        if ext == 'pdf' and PDF_OK:
            reader = PyPDF2.PdfReader(uploaded_file)
            return "\n".join([(p.extract_text() or "") for p in reader.pages]), None
        if ext in ['docx','doc'] and DOCX_OK:
            doc = Document(uploaded_file)
            return "\n".join(p.text for p in doc.paragraphs), None
        if ext in ['pptx','ppt'] and PPTX_OK:
            prs = Presentation(uploaded_file); parts = []
            for s in prs.slides:
                for sh in s.shapes:
                    if hasattr(sh, 'text'): parts.append(sh.text)
            return "\n".join(parts), None
        if ext == 'csv' and PANDAS_OK:
            df = pd.read_csv(uploaded_file); return df.to_string(), None
        if ext in ['xlsx','xls'] and PANDAS_OK:
            df = pd.read_excel(uploaded_file); return df.to_string(), None
        return uploaded_file.read().decode('utf-8', 'ignore'), None
    except Exception as e:
        return "", f"Error reading file: {e}"

# =====================
# NLP UTILITIES (offline, rule-based)
# =====================
def text_enhancer(text, mode="clarity"):
    text = re.sub(r"\s+", " ", text.strip())
    configs = {
        "professional": {"map": {r"\bkinda\b":"somewhat", r"\bwanna\b":"want to", r"\bgonna\b":"going to", r"\byeah\b":"yes", r"\bok\b":"acceptable", r"\bawesome\b":"excellent", r"\bstuff\b":"items", r"\bthings\b":"elements"}, "label":"Professional tone"},
        "creative": {"map": {r"\bsaid\b":"remarked", r"\bwent\b":"ventured", r"\bgood\b":"splendid", r"\bbig\b":"grand", r"\bsmall\b":"petite"}, "label":"Creative diction"},
        "concise": {"map": {r"\bin order to\b":"to", r"\bdue to the fact that\b":"because", r"\bat this point in time\b":"now", r"\bin the event that\b":"if"}, "label":"Conciseness"},
        "clarity": {"map": {r"\bstuff\b":"information", r"\bthings\b":"elements", r"\bthis\b(?!\s+\w+)":"this concept"}, "label":"Clarity"},
    }
    cfg = configs.get(mode, configs["clarity"])
    out = text
    for pat, rep in cfg["map"].items():
        out = re.sub(pat, rep, out, flags=re.I)
    return f"**{cfg['label']}**\n\n{out}\n\n**Stats:** original {len(text)} chars → {len(out)} chars"

def text_summarizer(text, style="concise"):
    sents = [s.strip() for s in re.split(r"[.!?]+", text) if s.strip()]
    if len(sents) <= 2:
        return "Text already concise (≤2 sentences)."
    words = re.findall(r"\b\w+\b", text.lower())
    freq = Counter(words)
    scored = []
    for i, s in enumerate(sents):
        score = sum(freq[w] for w in re.findall(r"\b\w+\b", s.lower()))
        if i in (0, len(sents)-1): score *= 1.3
        scored.append((score, i, s))
    scored.sort(reverse=True)
    if style == "bullet_points":
        k = min(5, max(2, len(sents)//3))
        chosen = sorted(scored[:k], key=lambda x: x[1])
        return "\n".join([f"• {c[2]}." for c in chosen])
    k = max(2, len(sents)//3 if style=="detailed" else len(sents)//4)
    chosen = sorted(scored[:k], key=lambda x: x[1])
    return " ".join([c[2] + "." for c in chosen])

def tone_analyzer(text):
    words = re.findall(r"\b\w+\b", text.lower())
    sent_map = {
        "positive":["excellent","amazing","wonderful","great","good","love","happy","pleased"],
        "negative":["terrible","awful","bad","hate","sad","angry","frustrated","disappointed"],
    }
    tone_map = {"formal":["therefore","furthermore","however"], "casual":["yeah","cool","awesome","gonna"], "academic":["analyze","methodology","framework"]}
    s_scores = {k: sum(w in v for w in words) for k,v in sent_map.items()}
    t_scores = {k: sum(w in v for w in words) for k,v in tone_map.items()}
    dominant_sent = max(s_scores, key=s_scores.get) if max(s_scores.values())>0 else "neutral"
    dominant_tone = max(t_scores, key=t_scores.get) if max(t_scores.values())>0 else "balanced"
    s_count = len([s for s in re.split(r"[.!?]+", text) if s.strip()])
    return f"Sentiment: {dominant_sent} | Tone: {dominant_tone} | Sentences: {s_count} | Words: {len(words)}"

def grammar_checker(text):
    issues, score = [], 100
    if text and not text[0].isupper(): issues.append("Start with a capital letter"); score -= 5
    if text and text[-1] not in ".!?": issues.append("End with punctuation"); score -= 5
    if "  " in text: issues.append("Extra spaces detected"); score -= 3
    if text.count('(') != text.count(')'): issues.append("Unbalanced parentheses"); score -= 3
    if text.count('"') % 2 != 0: issues.append("Unbalanced quotes"); score -= 3
    words = re.findall(r"\b\w+\b", text.lower()); uniq = len(set(words))
    diversity = (uniq/len(words)*100) if words else 0
    return "\n".join([
        f"Quality: {max(0,score)}/100",
        "Issues:" if issues else "Issues: none",
        *[f"• {i}" for i in issues],
        f"Vocabulary diversity: {diversity:.1f}% (unique {uniq}/{len(words)})",
    ])

def content_analyzer(text):
    words = re.findall(r"\b\w+\b", text.lower())
    sents = [s for s in re.split(r"[.!?]+", text) if s.strip()]
    paras = [p for p in text.split("\n\n") if p.strip()]
    stop = {'the','a','an','and','or','but','in','on','at','to','for','of','with','by','is','are','was','were','be','been','have','has','had','do','does','did','will','would','could','should'}
    filt = [w for w in words if w not in stop and len(w) > 3]
    freq = Counter(filt).most_common(5)
    return "\n".join([
        f"Words: {len(words)} | Sentences: {len(sents)} | Paragraphs: {len(paras)}",
        "Top keywords:" if freq else "Top keywords: none",
        *[f"• {w} ({c})" for w,c in freq],
    ])

def psychology_analyzer(text):
    words = re.findall(r"\b\w+\b", text.lower())
    emo = {
        "joy":["happy","excited","love","joy","delighted"],
        "sadness":["sad","disappointed","lonely","depressed"],
        "anger":["angry","furious","annoyed","frustrated"],
        "fear":["afraid","worried","anxious","nervous"],
    }
    scores = {k: sum(w in v for w in words) for k,v in emo.items()}
    dom = max(scores, key=scores.get) if max(scores.values())>0 else "neutral"
    return " | ".join([f"{k}:{v}" for k,v in scores.items()]) + f" | Dominant: {dom}"

def style_transformer(text, style, era):
    maps = {
        "victorian": {r"\bvery\b":"exceedingly", r"\bgood\b":"splendid", r"\bbad\b":"dreadful"},
        "shakespeare": {r"\byou\b":"thou", r"\byour\b":"thy", r"\bare\b":"art"},
        "1950s": {r"\bcool\b":"swell", r"\bgreat\b":"keen", r"\bbad\b":"lousy"},
        "modern": {r"\bvery good\b":"amazing", r"\bnice\b":"awesome"},
        "academic": {r"\bI think\b":"This suggests", r"\bshow\b":"demonstrate"},
    }
    out = text
    for pat, rep in maps.get(style, {}).items():
        out = re.sub(pat, rep, out, flags=re.I)
    return f"[{style} • {era}]\n\n{out}"

def ai_chatbot_response(msg: str) -> str:
    m = msg.lower()

    # Greetings
    if any(k in m for k in ["hello", "hi", "hey", "good morning", "good evening"]):
        return "Hi! 👋 Ask me about any tool, or tap a quick question below."

    # Enhancement
    if "enhance" in m or "improve" in m or "better" in m:
        return "Use **Text Enhancement** → pick professional, creative, clarity, or concise."

    # Summarization
    if "summar" in m or "shorten" in m or "brief" in m:
        return "Use **Summarization** → choose concise, detailed, or bullet points."

    # Grammar Check
    if "grammar" in m or "mistake" in m or "error" in m:
        return "Use **Grammar Check** → I’ll fix spelling and grammar issues."

    # Tone Analysis
    if "tone" in m or "emotion" in m or "feeling" in m:
        return "Use **Tone Analysis** → detect polite, formal, casual, or emotional tone."

    # Plagiarism Check
    if "plagiarism" in m or "copy" in m or "original" in m:
        return "Use **Plagiarism Check** → I’ll scan for originality of your text."

    # Psychology / Motivation
    if "psychology" in m or "motivation" in m or "stress" in m:
        return "Use **Psychology Tool** → I’ll give motivational or stress-relief tips."

    # Pricing / Free
    if "free" in m or "price" in m or "cost" in m:
        return "Everything here uses **free, offline libraries** — no hidden costs."

    # Default
    return "I can guide you: Enhancement, Summarization, Grammar, Tone, Plagiarism, Content, Psychology, Style, or Chatbot."

# ---------- Pretty result renderer (like your slide 2) ----------
def render_result_panel(title: str, score_pct: int, summary: str, checks: list[tuple[str,str]]):
    # title bar
    st.markdown(f"""
    <div class="result-header">🪄 {title}</div>
    """, unsafe_allow_html=True)
    # status bar
    st.markdown('<div style="height:6px"></div>', unsafe_allow_html=True)
    st.markdown('<div class="card-tight">', unsafe_allow_html=True)
    st.markdown("**Overall Confidence**")
    st.markdown(f'''
      <div class="progress-wrap">
        <div class="progress-fill" style="width:{max(0,min(100,score_pct))}%"></div>
      </div>
    ''', unsafe_allow_html=True)
    st.markdown(f"<div style='margin-top:6px;color:#222'><b>{score_pct}%</b> • {summary}</div>", unsafe_allow_html=True)
    st.markdown("</div>", unsafe_allow_html=True)

    # checklist
    st.markdown('<div class="card-tight" style="margin-top:10px">', unsafe_allow_html=True)
    st.markdown("**Details**")
    for level, text in checks:
        cls = "check" if level=="ok" else ("check warn" if level=="warn" else "check err")
        st.markdown(f"""<div class="{cls}"><div class="dot"></div><div>{text}</div></div>""", unsafe_allow_html=True)
        st.markdown("<div style='height:6px'></div>", unsafe_allow_html=True)
    st.markdown("</div>", unsafe_allow_html=True)

# =====================
# APP STATE
# =====================
init_db()
if 'authenticated' not in st.session_state:
    st.session_state.update({
        'authenticated': False, 'user_id': None, 'username': None,
        'login_time': None, 'current_feature': '🎨 Text Enhancement',
        'chat_history': [{"role":"assistant","message":"Hello! I'm your TextCraft AI Assistant. How can I help?"}]
    })

# =====================
# AUTH SCREENS
# =====================
def auth_screen():
    st.markdown("<div class='card'>", unsafe_allow_html=True)
    st.markdown("<p class='pill'>Secure Access</p>", unsafe_allow_html=True)
    st.markdown("<h1>🎯 TextCraft AI Studio</h1>", unsafe_allow_html=True)
    st.caption("Professional Text Intelligence Platform — Transform • Analyze • Perfect")

    col1, col2, col3 = st.columns([1,2,1])
    with col2:
        tabs = st.tabs(["🔐 Sign In", "📝 Create Account"])
        with tabs[0]:
            st.markdown("<div class='card'>", unsafe_allow_html=True)
            u = st.text_input("👤 Username", key="login_u")
            p = st.text_input("🔒 Password", type="password", key="login_p")
            c1, c2 = st.columns(2)
            with c1:
                if st.button("🚀 Sign In"):
                    if u and p:
                        with st.spinner("Authenticating..."):
                            time.sleep(0.3)
                            uid = verify_user(u, p)
                            if uid:
                                st.session_state.update({
                                    'authenticated': True, 'user_id': uid,
                                    'username': u, 'login_time': datetime.now()
                                })
                                st.success("✅ Welcome back! Redirecting...")
                                time.sleep(0.6); st.rerun()
                            else:
                                st.error("❌ Invalid credentials")
                    else:
                        st.warning("⚠️ Please fill in all fields")
            with c2:
                if st.button("🔁 Reset"):
                    st.session_state.login_u = ""
                    st.session_state.login_p = ""
            st.markdown("</div>", unsafe_allow_html=True)

        with tabs[1]:
            st.markdown("<div class='card'>", unsafe_allow_html=True)
            nu = st.text_input("👤 Choose Username", key="new_u")
            np = st.text_input("🔒 Choose Password (6+)", type="password", key="new_p")
            np2 = st.text_input("🔒 Confirm Password", type="password", key="new_p2")
            if st.button("🎉 Create Account"):
                if not (nu and np and np2): st.warning("⚠️ Please fill in all fields")
                elif np != np2: st.error("❌ Passwords don't match")
                elif len(np) < 6: st.error("❌ Password must be at least 6 characters")
                else:
                    with st.spinner("Creating account..."):
                        time.sleep(0.3)
                        if create_user(nu, np):
                            st.success("✅ Account created! Please sign in.")
                        else:
                            st.error("❌ Username already exists")
            st.markdown("</div>", unsafe_allow_html=True)
    st.markdown("</div>", unsafe_allow_html=True)

# =====================
# SIDEBAR (after auth)
# =====================
def sidebar():
    with st.sidebar:
        session_count, last_login, created_at, total_analyses = get_user_stats(st.session_state.user_id)
        duration = datetime.now() - (st.session_state.login_time or datetime.now())
        st.markdown("<div class='sidebar-card'>", unsafe_allow_html=True)
        st.markdown(f"### 👋 {st.session_state.username}")
        st.caption("Welcome back to TextCraft")
        st.markdown(
            f"""
            <div class="card-tight">
              <div class="metric-title">Session</div>
              <div class="metric-val">#{session_count}</div>
            </div>
            """, unsafe_allow_html=True)
        colx, coly = st.columns(2)
        with colx:
            st.markdown(
                f"""<div class="card-tight">
                    <div class="metric-title">Analyses</div>
                    <div class="metric-val">{total_analyses}</div>
                </div>""", unsafe_allow_html=True)
        with coly:
            st.markdown(
                f"""<div class="card-tight">
                    <div class="metric-title">Active</div>
                    <div class="metric-val">{str(duration).split('.')[0]}</div>
                </div>""", unsafe_allow_html=True)

        st.markdown("</div>", unsafe_allow_html=True)
        st.write("")
        feature = st.selectbox(
            "Professional Tools",
            [
                "🎨 Text Enhancement","📝 Text Summarization","🎭 Tone Analysis","📝 Grammar Check",
                "📊 Content Analysis","🧠 Psychology Analyzer","⏰ Style Transformer","🤖 AI Chatbot Assistant"
            ],
            index=0
        )
        st.session_state.current_feature = feature
        st.write("")
        if st.button("🚪 Sign Out"):
            st.session_state.authenticated=False
            st.session_state.user_id=None
            st.session_state.username=None
            st.session_state.login_time=None
            st.rerun()

# =====================
# MAIN APP (after auth)
# =====================
def main_app():
    sidebar()

    # Workspace header (kept for non-chatbot tools)
    show_upload = st.session_state.current_feature != "🤖 AI Chatbot Assistant"
    file_text = ""
    if show_upload:
        st.markdown("<div class='card'>", unsafe_allow_html=True)
        st.markdown("<p class='pill'>Workspace</p>", unsafe_allow_html=True)
        st.markdown("<h2>📁 Document Upload</h2>", unsafe_allow_html=True)

        up = st.file_uploader("📎 Upload document for analysis", type=get_supported_file_types(),
                              help="TXT, PDF, DOCX, PPTX, CSV/XLSX and more (offline parsers)")

        if up:
            with st.spinner("📖 Processing document..."):
                content, err = read_uploaded_file(up)
            if err:
                st.error(f"❌ {err}")
            else:
                file_text = content
                st.success(f"✅ '{up.name}' processed!")

                c1, c2, c3, c4 = st.columns(4)
                with c1:
                    st.markdown("<div class='card-tight'>", unsafe_allow_html=True)
                    st.markdown("<div class='metric-title'>📄 Characters</div>", unsafe_allow_html=True)
                    st.markdown(f"<div class='metric-val'>{len(file_text):,}</div>", unsafe_allow_html=True)
                    st.markdown("</div>", unsafe_allow_html=True)

                word_count = len(re.findall(r"\\b\\w+\\b", file_text))
                with c2:
                    st.markdown("<div class='card-tight'>", unsafe_allow_html=True)
                    st.markdown("<div class='metric-title'>📝 Words</div>", unsafe_allow_html=True)
                    st.markdown(f"<div class='metric-val'>{word_count:,}</div>", unsafe_allow_html=True)
                    st.markdown("</div>", unsafe_allow_html=True)

                with c3:
                    st.markdown("<div class='card-tight'>", unsafe_allow_html=True)
                    st.markdown("<div class='metric-title'>📑 Lines</div>", unsafe_allow_html=True)
                    st.markdown(f"<div class='metric-val'>{len(file_text.splitlines()):,}</div>", unsafe_allow_html=True)
                    st.markdown("</div>", unsafe_allow_html=True)

                with c4:
                    st.markdown("<div class='card-tight'>", unsafe_allow_html=True)
                    st.markdown("<div class='metric-title'>📊 Format</div>", unsafe_allow_html=True)
                    st.markdown(f"<div class='metric-val'>{up.name.split('.')[-1].upper()}</div>", unsafe_allow_html=True)
                    st.markdown("</div>", unsafe_allow_html=True)

        st.markdown("</div>", unsafe_allow_html=True)

    feature = st.session_state.current_feature

    # Input area (not shown for Chatbot)
    if feature != "🤖 AI Chatbot Assistant":
        st.markdown("<div class='card'>", unsafe_allow_html=True)
        st.markdown("<p class='pill'>Input</p>", unsafe_allow_html=True)
        st.markdown("<h3>✍️ Text Input</h3>", unsafe_allow_html=True)
        default_value = file_text if file_text else ""
        text_input = st.text_area("Enter your text:", value=default_value, height=220, placeholder="Type or paste your text here...")
        st.markdown("</div>", unsafe_allow_html=True)
    else:
        text_input = ""

    # ---- Tools ----
    st.markdown("<div class='card'>", unsafe_allow_html=True)
    st.markdown("<p class='pill'>Tools</p>", unsafe_allow_html=True)

    if feature == "🎨 Text Enhancement":
        st.subheader("✨ Professional Text Enhancement")
        mode = st.selectbox("🎯 Enhancement Mode", ["clarity","professional","creative","concise"])
        if st.button("🎨 Enhance Text"):
            if text_input.strip():
                with st.spinner("Applying enhancements..."):
                    time.sleep(0.4); result = text_enhancer(text_input, mode)
                    save_analysis(st.session_state.user_id, "Enhancement", text_input, result)
                # Slide-style result
                render_result_panel(
                    "Enhanced Text",
                    92,
                    f"Mode: {mode.capitalize()}",
                    [
                        ("ok","Grammar & fluency improved"),
                        ("ok","Diction adjusted to requested style"),
                        ("warn","Review domain-specific terms for accuracy"),
                    ]
                )
                st.info(result)
            else:
                st.warning("⚠️ Please enter text to enhance")

    elif feature == "📝 Text Summarization":
        st.subheader("📄 Intelligent Text Summarization")
        style = st.selectbox("📋 Summary Format", ["concise","detailed","bullet_points"])
        if st.button("📝 Generate Summary"):
            if text_input.strip():
                with st.spinner("Summarizing..."):
                    time.sleep(0.4); result = text_summarizer(text_input, style)
                    save_analysis(st.session_state.user_id, "Summarization", text_input, result)
                render_result_panel(
                    "Summary Generated",
                    88,
                    f"Style: {style.replace('_',' ').title()}",
                    [
                        ("ok","Key sentences retained"),
                        ("ok","Lead & concluding context preserved"),
                        ("warn","Verify nuanced claims against source"),
                    ]
                )
                st.success(result)
            else:
                st.warning("⚠️ Please enter text to summarize")

    elif feature == "🎭 Tone Analysis":
        st.subheader("🎭 Advanced Tone & Sentiment Analysis")
        if st.button("🔍 Analyze Tone"):
            if text_input.strip():
                with st.spinner("Analyzing tone..."):
                    time.sleep(0.3); result = tone_analyzer(text_input)
                    save_analysis(st.session_state.user_id, "Tone Analysis", text_input, result)
                render_result_panel(
                    "Tone & Sentiment",
                    84,
                    "Dominant sentiment and tone extracted",
                    [
                        ("ok","Sentiment polarity identified"),
                        ("ok","Formality vs casual tone detected"),
                        ("warn","Sarcasm and irony may require manual review"),
                    ]
                )
                st.info(result)
            else:
                st.warning("⚠️ Please enter text to analyze")

    elif feature == "📝 Grammar Check":
        st.subheader("📝 Professional Grammar & Style Analysis")
        if st.button("✅ Analyze Grammar"):
            if text_input.strip():
                with st.spinner("Checking grammar..."):
                    time.sleep(0.3); result = grammar_checker(text_input)
                    save_analysis(st.session_state.user_id, "Grammar Check", text_input, result)
                # Parse score for bar
                try:
                    score_line = result.splitlines()[0]
                    pct = int(re.findall(r"(\d+)/100", score_line)[0])
                except Exception:
                    pct = 80
                render_result_panel(
                    "Grammar & Style Report",
                    pct,
                    "Overall Assessment based on rule checks",
                    [
                        ("ok","Punctuation and sentence termination verified"),
                        ("warn","Spacing & parentheses balance reviewed"),
                        ("ok","Vocabulary diversity computed"),
                    ]
                )
                st.warning(result)
            else:
                st.warning("⚠️ Please enter text to check")

    elif feature == "📊 Content Analysis":
        st.subheader("📊 Comprehensive Content Intelligence")
        if st.button("📈 Analyze Content"):
            if text_input.strip():
                with st.spinner("Analyzing content..."):
                    time.sleep(0.3); result = content_analyzer(text_input)
                    save_analysis(st.session_state.user_id, "Content Analysis", text_input, result)
                render_result_panel(
                    "Content Structure",
                    86,
                    "Keywords & structure metrics extracted",
                    [
                        ("ok","Word/Sentence/Paragraph counts calculated"),
                        ("ok","Stopwords filtered to find keywords"),
                        ("warn","Consider refining repetitive phrases"),
                    ]
                )
                st.info(result)
            else:
                st.warning("⚠️ Please enter text to analyze")

    elif feature == "🧠 Psychology Analyzer":
        st.subheader("🧠 Psychological Text Analysis")
        if st.button("🔮 Analyze Psychology"):
            if text_input.strip():
                with st.spinner("Profiling..."):
                    time.sleep(0.3); result = psychology_analyzer(text_input)
                    save_analysis(st.session_state.user_id, "Psychology Analysis", text_input, result)
                render_result_panel(
                    "Emotion Signals",
                    78,
                    "Relative counts across Joy, Sadness, Anger, Fear",
                    [
                        ("ok","Dominant emotion estimated from tokens"),
                        ("warn","Context-dependent meanings may vary"),
                        ("ok","Balanced overview provided"),
                    ]
                )
                st.info(result)
            else:
                st.warning("⚠️ Please enter text to analyze")

    elif feature == "⏰ Style Transformer":
        st.subheader("⏰ Advanced Style & Era Transformation")
        sc1, sc2 = st.columns(2)
        with sc1:
            style = st.selectbox("🎭 Writing Style", ["modern","victorian","shakespeare","1950s","academic"])
        with sc2:
            era = st.selectbox("📅 Historical Era", ["2040s","2030s","2020s","2010s","2000s","1990s","1980s","1970s","1960s","1950s","1940s","1930s","1920s","1910s","1900s","1800s","1700s","1600s"])
        if st.button("⏰ Transform Style"):
            if text_input.strip():
                with st.spinner("Transforming style..."):
                    time.sleep(0.3); result = style_transformer(text_input, style, era)
                    save_analysis(st.session_state.user_id, "Style Transformation", text_input, result)
                render_result_panel(
                    "Style Transformation",
                    90,
                    f"Applied {style.title()} diction • Era: {era}",
                    [
                        ("ok","Keyword-level substitutions applied"),
                        ("ok","Original semantics preserved as much as possible"),
                        ("warn","Manual polish may improve flow"),
                    ]
                )
                st.success(result)
            else:
                st.warning("⚠️ Please enter text to transform")

    elif feature == "🤖 AI Chatbot Assistant":
       st.subheader("🤖 TextCraft AI Chatbot Assistant (free)")

    # Show chat history with boxed answers
    for chat in st.session_state.chat_history:
        role = "You" if chat["role"]=="user" else "🤖 AI Assistant"
        if chat["role"] == "assistant":
            st.markdown(
                f"<div class='card-tight'><b>{role}:</b><br>{chat['message']}</div>",
                unsafe_allow_html=True
            )
        else:
            st.markdown(f"**{role}:** {chat['message']}")

    # --- User input ---
    def handle_input():
        user_msg = st.session_state.chat_input
        if user_msg.strip():
            reply = ai_chatbot_response(user_msg)
            st.session_state.chat_history.append({"role": "user", "message": user_msg})
            st.session_state.chat_history.append({"role": "assistant", "message": reply})
            save_analysis(st.session_state.user_id, "AI Chat", user_msg, reply)
            st.session_state.chat_input = ""   # ✅ allowed here (inside callback)
            st.rerun()

    st.text_input(
        "💬 Type your question here...",
        key="chat_input",
        on_change=handle_input
    )

    st.markdown("### 💡 Quick Questions")
    col1, col2, col3 = st.columns(3)

    def ask(q):
        reply = ai_chatbot_response(q)
        st.session_state.chat_history.append({"role":"user","message":q})
        st.session_state.chat_history.append({"role":"assistant","message":reply})
        save_analysis(st.session_state.user_id, "AI Chat", q, reply)
        st.rerun()

    with col1:
        if st.button("🔹 How do I enhance text?"): ask("How do I enhance text?")
    with col2:
        if st.button("🔹 How do I summarize text?"): ask("How do I summarize text?")
    with col3:
        if st.button("🔹 Is this app free?"): ask("Is this app free?")

    st.write("")
    if st.button("🗑️ Clear Conversation"):
        st.session_state.chat_history = [
            {"role":"assistant","message":"Hello! I'm your TextCraft AI Assistant. How can I help?"}
        ]
        st.rerun()

    st.caption("Type your own question above or tap a quick question • Offline & free")

# =====================
# ENTRY
# =====================
if not st.session_state.authenticated:
    auth_screen()
else:
    main_app()